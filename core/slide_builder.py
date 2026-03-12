"""
S4 Graficador — Slide Builder (v2 — python-pptx)
Usa python-pptx para duplicar slides (mantiene slide master/layout válidos),
y cirugía XML solo para inyectar datos en charts.
"""

from io import BytesIO
from copy import deepcopy
from lxml import etree
from typing import List, Dict, Optional
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI

from .models import Question, ParseResult, AGORA_STYLE
from .chart_injector import inject_chart_data, clean_chart_xml


# ──────────────────────────────────────────────────────────────────────────────
# TEMPLATE SLIDE INDEX MAPPING (0-based, en Template_clean.pptx)
# ──────────────────────────────────────────────────────────────────────────────

TEMPLATE_SLIDE_IDX = {
    "FREC_MULTIPLE": 0,        # slide 0: 1 chart (percentStacked col)
    "FREC_MULTI_GRAFICOS": 2,  # slide 2: 3 charts (clustered col)
    "FREC_SIMPLE": 4,          # slide 4: 1 chart (clustered col)
    "APERTURA_SIMPLE": 5,      # slide 5: 1 chart (percentStacked bar)
}

# Namespaces for XML surgery
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_CHART = "http://schemas.openxmlformats.org/drawingml/2006/chart"

# Path al template limpio (sin image4.jpeg corrupta)
_CLEAN_TEMPLATE = Path(__file__).resolve().parent.parent / "templates" / "Template_clean.pptx"

# Global chart counter para nombres únicos
_chart_counter = 0


def _next_chart_num():
    global _chart_counter
    _chart_counter += 1
    return _chart_counter


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE ORDERING LOGIC (sin cambios respecto a v1)
# ──────────────────────────────────────────────────────────────────────────────

def _compute_slide_order(questions: List[Question]) -> List[dict]:
    slides = []
    cap_order = []
    seen_caps = set()
    for q in questions:
        if q.capitulo not in seen_caps:
            cap_order.append(q.capitulo)
            seen_caps.add(q.capitulo)

    app_chapters = {}
    for q in questions:
        if q.cap_app:
            app_chapters.setdefault(q.cap_app, []).append(q)

    for cap in cap_order:
        cap_questions = [q for q in questions if q.capitulo == cap]
        grupos = {}
        for q in cap_questions:
            grupos.setdefault(q.grupo_frec, []).append(q)

        for grupo_id in sorted(grupos.keys()):
            grupo_qs = grupos[grupo_id]
            tipo = grupo_qs[0].tipo_slide

            if tipo == "FREC_SIMPLE":
                for q in grupo_qs:
                    slides.append({
                        "type": "FREC", "chart_type": "FREC_SIMPLE",
                        "question": q, "questions_group": None,
                        "titulo_slide": f"{q.titulo_frec} | {q.titulo_app}",
                        "titulo_chart": q.titulo_app, "capitulo": cap,
                    })
                    if not q.cap_app:
                        slides.append({
                            "type": "APERTURA", "chart_type": "APERTURA_SIMPLE",
                            "question": q, "questions_group": None,
                            "titulo_slide": q.titulo_app,
                            "titulo_chart": q.titulo_app, "capitulo": cap,
                        })

            elif tipo == "FREC_MULTIPLE":
                slides.append({
                    "type": "FREC", "chart_type": "FREC_MULTIPLE",
                    "question": None, "questions_group": grupo_qs,
                    "titulo_slide": f"{grupo_qs[0].titulo_frec} | {grupo_qs[0].titulo_app}",
                    "titulo_chart": grupo_qs[0].titulo_app, "capitulo": cap,
                })
                for q in grupo_qs:
                    if not q.cap_app:
                        slides.append({
                            "type": "APERTURA", "chart_type": "APERTURA_SIMPLE",
                            "question": q, "questions_group": None,
                            "titulo_slide": q.titulo_app,
                            "titulo_chart": q.titulo_app, "capitulo": cap,
                        })

            elif tipo == "FREC_MULTI_GRAFICOS":
                slides.append({
                    "type": "FREC", "chart_type": "FREC_MULTI_GRAFICOS",
                    "question": None, "questions_group": grupo_qs,
                    "titulo_slide": f"{grupo_qs[0].titulo_frec} | {grupo_qs[0].titulo_app}",
                    "titulo_chart": grupo_qs[0].titulo_app, "capitulo": cap,
                })
                for q in grupo_qs:
                    if not q.cap_app:
                        slides.append({
                            "type": "APERTURA", "chart_type": "APERTURA_SIMPLE",
                            "question": q, "questions_group": None,
                            "titulo_slide": q.titulo_app,
                            "titulo_chart": q.titulo_app, "capitulo": cap,
                        })

    for cap_app in sorted(app_chapters.keys()):
        for q in app_chapters[cap_app]:
            slides.append({
                "type": "APERTURA", "chart_type": "APERTURA_SIMPLE",
                "question": q, "questions_group": None,
                "titulo_slide": q.titulo_app,
                "titulo_chart": q.titulo_app, "capitulo": cap_app,
            })

    return slides


# ──────────────────────────────────────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────────────────────────────────────

def _duplicate_slide(prs, src_slide):
    """
    Duplica una slide en la presentación usando python-pptx.
    Crea chart parts con nombres únicos para evitar duplicados en el ZIP.
    Retorna (new_slide, chart_rids_map) donde chart_rids_map es {rId: new_chart_part}.
    """
    layout = src_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # Collect source chart rels BEFORE copying spTree
    src_chart_rels = {}  # {rId: chart_part}
    for rel in src_slide.part.rels.values():
        if "chart" in rel.reltype:
            src_chart_rels[rel.rId] = (rel.reltype, rel.target_part)

    # Copy spTree from source to new slide
    src_cSld = src_slide._element.find(f"{{{NS_P}}}cSld")
    new_cSld = new_slide._element.find(f"{{{NS_P}}}cSld")

    if src_cSld is not None and new_cSld is not None:
        # Remove new slide's spTree (auto-generated empty one)
        new_spTree = new_cSld.find(f"{{{NS_P}}}spTree")
        if new_spTree is not None:
            new_cSld.remove(new_spTree)
        # Copy source's spTree
        src_spTree = src_cSld.find(f"{{{NS_P}}}spTree")
        if src_spTree is not None:
            new_spTree = deepcopy(src_spTree)
            new_cSld.append(new_spTree)

    # Create new chart parts with unique names and add rels
    # Map: old rId -> new rId
    rid_map = {}
    new_chart_parts = {}  # {new_rId: new_part}

    for old_rId, (reltype, src_chart_part) in src_chart_rels.items():
        chart_num = _next_chart_num()
        new_partname = PackURI(f'/ppt/charts/chart{chart_num}.xml')

        # Create new Part with unique name, copy blob from source
        new_chart = Part(
            new_partname,
            src_chart_part.content_type,
            prs.part.package,
            src_chart_part.blob,
        )

        # DON'T copy chart sub-rels (themeOverride, style, colors, oleObject)
        # chart_injector strips externalData/clrMapOvr refs from the XML

        # Add relationship from new slide to new chart part
        new_rId = new_slide.part.relate_to(new_chart, reltype)
        rid_map[old_rId] = new_rId
        new_chart_parts[new_rId] = new_chart

    # Update r:id references in the copied spTree to point to new chart rels
    if rid_map:
        _update_chart_rids(new_slide._element, rid_map)

    return new_slide, new_chart_parts


def _update_chart_rids(slide_element, rid_map):
    """
    Updates r:id references in graphicFrame elements to use new rIds.
    """
    # Find all c:chart elements (inside graphicData)
    for chart_ref in slide_element.iter(f"{{{NS_CHART}}}chart"):
        old_rid = chart_ref.get(f"{{{NS_R}}}id")
        if old_rid and old_rid in rid_map:
            chart_ref.set(f"{{{NS_R}}}id", rid_map[old_rid])


def _get_chart_parts_from_map(chart_parts_map):
    """Returns list of chart parts from the map, in rId order."""
    return [chart_parts_map[k] for k in sorted(chart_parts_map.keys())]


def _clean_shapes(slide):
    """
    Limpia shapes del template: textboxes sueltos, grupos, conectores.
    Mantiene: placeholders (title, ftr, sldNum) y graphicFrames (charts).
    """
    spTree = slide._element.find(f".//{{{NS_P}}}cSld/{{{NS_P}}}spTree")
    if spTree is None:
        return

    to_remove = []
    for child in list(spTree):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

        if tag == "sp":
            nvSpPr = child.find(f"{{{NS_P}}}nvSpPr")
            if nvSpPr is not None:
                nvPr = nvSpPr.find(f"{{{NS_P}}}nvPr")
                if nvPr is not None:
                    ph = nvPr.find(f"{{{NS_P}}}ph")
                    if ph is not None:
                        ph_type = ph.get("type", "")
                        if ph_type in ("title", "ctrTitle", "ftr", "sldNum"):
                            continue
            to_remove.append(child)
        elif tag in ("grpSp", "cxnSp"):
            to_remove.append(child)

    for el in to_remove:
        spTree.remove(el)


def _set_title(slide, titulo):
    """Setea el título del slide usando el placeholder."""
    if not titulo:
        return
    for shape in slide.placeholders:
        if shape.placeholder_format.type is not None:
            idx = shape.placeholder_format.idx
            if idx == 0:  # Title placeholder
                shape.text = titulo
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = AGORA_STYLE["slide_title_font"]
                        run.font.size = Emu(AGORA_STYLE["slide_title_size"] * 12700 // 100)
                        run.font.bold = AGORA_STYLE["slide_title_bold"]
                return


def _set_footer(slide, footer_text):
    """Setea el pie de página del slide."""
    if not footer_text:
        return
    for shape in slide.placeholders:
        if hasattr(shape.placeholder_format, 'type'):
            ph_el = shape._element.find(f".//{{{NS_P}}}nvSpPr/{{{NS_P}}}nvPr/{{{NS_P}}}ph")
            if ph_el is not None and ph_el.get("type") == "ftr":
                shape.text = footer_text
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = AGORA_STYLE["footer_font"]
                        run.font.size = Emu(AGORA_STYLE["footer_size"] * 12700 // 100)
                        run.font.bold = AGORA_STYLE["footer_bold"]
                return


def _expand_chart_for_apertura(slide):
    """Expande el graphicFrame del chart para aperturas."""
    for gf in slide._element.iter(f"{{{NS_P}}}graphicFrame"):
        xfrm = gf.find(f".//{{{NS_P}}}xfrm")
        if xfrm is None:
            xfrm = gf.find(f".//{{{NS_A}}}xfrm")
        if xfrm is None:
            continue
        off = xfrm.find(f"{{{NS_A}}}off")
        ext = xfrm.find(f"{{{NS_A}}}ext")
        if off is not None and ext is not None:
            off.set("x", "500000")
            off.set("y", "2100000")
            ext.set("cx", "23200000")
            ext.set("cy", "9800000")


# ──────────────────────────────────────────────────────────────────────────────
# PUBLIC API
# ──────────────────────────────────────────────────────────────────────────────

def build_pptx(template_bytes: bytes, questions: List[Question],
               segment_groups=None, pie_pagina: str = "") -> bytes:
    """Genera un PPTX con gráficos a partir del template y las preguntas."""

    global _chart_counter

    # Abrir template limpio con python-pptx
    prs = Presentation(str(_CLEAN_TEMPLATE))

    # Determinar el máximo chart number existente en el template
    max_chart = 0
    for part in prs.part.package.iter_parts():
        pn = str(part.partname)
        if '/charts/chart' in pn and pn.endswith('.xml') and 'style' not in pn and 'colors' not in pn:
            try:
                num = int(pn.split('chart')[-1].replace('.xml', ''))
                max_chart = max(max_chart, num)
            except ValueError:
                pass
    _chart_counter = max_chart  # start from max + 1

    # Guardar slides del template como fuentes para duplicar
    template_slides = list(prs.slides)

    slide_order = _compute_slide_order(questions)

    # Crear todas las slides nuevas duplicando del template
    new_slides = []
    for spec in slide_order:
        chart_type = spec["chart_type"]
        src_idx = TEMPLATE_SLIDE_IDX.get(chart_type)
        if src_idx is None or src_idx >= len(template_slides):
            continue
        src_slide = template_slides[src_idx]
        new_slide, chart_parts = _duplicate_slide(prs, src_slide)
        new_slides.append((new_slide, spec, chart_parts))

    # Eliminar las slides originales del template (con drop_rel para limpiar parts)
    from pptx.oxml.ns import qn
    xml_slides = prs.slides._sldIdLst
    original_ids = list(xml_slides)[:len(template_slides)]
    for sldId in original_ids:
        rId = sldId.get(qn('r:id'))
        prs.part.drop_rel(rId)
        xml_slides.remove(sldId)

    # Limpiar XML de TODOS los chart parts (elimina r:id huérfanos, mc:AlternateContent)
    for _, _, chart_parts in new_slides:
        for part in chart_parts.values():
            part._blob = clean_chart_xml(part.blob)

    # Ahora procesar cada slide: limpiar, inyectar datos, título, footer
    for new_slide, spec, chart_parts in new_slides:
        chart_type = spec["chart_type"]

        # Limpiar shapes sobrantes
        _clean_shapes(new_slide)

        # Get chart parts in order
        parts_list = _get_chart_parts_from_map(chart_parts)

        if chart_type == "FREC_MULTI_GRAFICOS" and spec.get("questions_group"):
            grupo = spec["questions_group"]
            for i, q in enumerate(grupo):
                if i < len(parts_list):
                    part = parts_list[i]
                    new_xml = inject_chart_data(
                        chart_xml_bytes=part.blob,
                        chart_type="FREC_MULTI_GRAFICOS",
                        question=q,
                        title=q.titulo_app,
                    )
                    part._blob = new_xml
        else:
            # Single chart
            question = spec.get("question")
            questions_group = spec.get("questions_group")
            titulo_chart = spec.get("titulo_chart", "")
            orientacion = "V"
            if question:
                orientacion = getattr(question, "orientacion", "V")
            elif questions_group:
                orientacion = getattr(questions_group[0], "orientacion", "V")

            if parts_list:
                part = parts_list[0]
                new_xml = inject_chart_data(
                    chart_xml_bytes=part.blob,
                    chart_type=chart_type,
                    question=question,
                    questions_group=questions_group,
                    title=titulo_chart,
                    orientacion=orientacion,
                    segment_groups=segment_groups,
                )
                part._blob = new_xml

        # Expandir chart para aperturas
        if chart_type == "APERTURA_SIMPLE":
            _expand_chart_for_apertura(new_slide)

        # Título y footer
        _set_title(new_slide, spec.get("titulo_slide", ""))
        _set_footer(new_slide, pie_pagina)

    # Guardar a bytes
    output = BytesIO()
    prs.save(output)
    return output.getvalue()
